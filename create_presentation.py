from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper to add a slide with title and content
    def add_slide(title, content_points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(20)

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Digital Twin Project"
    subtitle.text = "Advanced LLM Orchestration & RAG Platform\nFrontend: React | Backend: Django"

    # 2. Project Overview
    add_slide("Project Overview", [
        "A comprehensive platform for interacting with, comparing, and evaluating multiple Large Language Models (LLMs).",
        "Key capabilities:",
        "- Multi-Model Chat (OpenAI, Mistral, Gemini)",
        "- RAG (Retrieval-Augmented Generation) for document context",
        "- Voice Mode for hands-free interaction",
        "- 'The Arena' for side-by-side model comparison"
    ])

    # 3. Core Features: Chat & RAG
    add_slide("Core Features: Chat & RAG", [
        "Unified Chat Interface: Switch seamlessly between different LLM providers.",
        "Document Integration (RAG):",
        "- Upload PDF, DOCX, TXT files",
        "- Chat with documents using context-aware retrieval",
        "- Persistent global document state"
    ])

    # 4. Core Features: Voice Mode
    add_slide("Core Features: Voice Chat", [
        "Full 'Speech-to-Speech' experience.",
        "Continuous Conversation Mode:",
        "- Auto-restarts microphone after AI speaks",
        "- Hands-free 'Phone Call' flow",
        "Minimalist UI with wave visualizations",
        "Powered by Web Speech API & Backend TTS"
    ])

    # 5. Core Features: The Arena
    add_slide("Core Features: The Arena", [
        "Model Battleground:",
        "- Compare responses from two models side-by-side",
        "- Blind evaluation or explicit judging",
        "Analytics & Leaderboard",
        "Automated AI Judge Service"
    ])

    # 6. Technical Architecture
    add_slide("Technical Architecture", [
        "Frontend: React, TailwindCSS, Vite",
        "- Context-based State Management",
        "- Real-time Audio Visualizers",
        "Backend: Django REST Framework",
        "- Services: LLMOrchestrator, JudgeService, AnalysisService",
        "- Vector Store: ChromaDB / FAISS for RAG",
        "LLM Integration: LangChain / Direct API"
    ])

    # 7. Future Roadmap
    add_slide("Future Roadmap", [
        "Agentic Workflows (Autonomous execution)",
        "Expanded Model Support (Llama 3, Claude 3)",
        "Mobile App Version",
        "Advanced Analytics Dashboard",
        "Multi-modal inputs (Image/Video analysis)"
    ])

    prs.save('Digital_Twin_Project_Presentation.pptx')
    print("Presentation saved as 'Digital_Twin_Project_Presentation.pptx'")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Error: python-pptx library not found. Please run 'pip install python-pptx'")
