import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import A4


class FileGenerator:
    # ---------------------------------------
    # DOCX GENERATION
    # ---------------------------------------

    def generate_docx(self, content: str, output_path="output.docx"):
        """
        Create a DOCX file from plain text.
        """
        doc = Document()

        for line in content.split("\n"):
            doc.add_paragraph(line)

        doc.save(output_path)
        return output_path

    # ---------------------------------------
    # PDF GENERATION
    # ---------------------------------------

    def generate_pdf(self, content: str, output_path="output.pdf"):
        """
        Create a PDF file using reportlab.
        """
        doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()
        flowables = []

        for line in content.split("\n"):
            flowables.append(Paragraph(line, styles["Normal"]))

        doc.build(flowables)
        return output_path

    # ---------------------------------------
    # PPTX GENERATION
    # ---------------------------------------

    def generate_pptx(self, content: str, output_path="output.pptx"):
        """
        Convert content into a PowerPoint file,
        with one bullet point per line.
        """
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title + body slide

        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "Generated Presentation"
        tf = body.text_frame

        lines = [line.strip() for line in content.split("\n") if line.strip()]

        if lines:
            tf.text = lines[0]
            for line in lines[1:]:
                p = tf.add_paragraph()
                p.text = line
                p.level = 1

        prs.save(output_path)
        return output_path


# ------------------------------
# Singleton instance for import
# ------------------------------

file_generator = FileGenerator()
