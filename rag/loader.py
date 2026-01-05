# rag/loader.py

import os
from typing import List
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    UnstructuredHTMLLoader,
)
from docx import Document
from pptx import Presentation


def load_pdf(path: str) -> List[dict]:
    docs = PyPDFLoader(path).load()
    return [{"content": d.page_content or "", "metadata": {"source": os.path.basename(path)}} for d in docs]


def load_txt(path: str) -> List[dict]:
    docs = TextLoader(path).load()
    return [{"content": d.page_content or "", "metadata": {"source": os.path.basename(path)}} for d in docs]


def load_docx(path: str) -> List[dict]:
    doc = Document(path)
    full_text = []
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
            
        style_name = para.style.name.lower()
        
        if 'heading 1' in style_name:
            full_text.append(f"# {text}")
        elif 'heading 2' in style_name:
            full_text.append(f"## {text}")
        elif 'heading 3' in style_name:
            full_text.append(f"### {text}")
        elif 'heading 4' in style_name:
            full_text.append(f"#### {text}")
        elif 'bullet' in style_name or 'list' in style_name:
            full_text.append(f"- {text}")
        else:
            full_text.append(text)
            
    content = "\n\n".join(full_text)
    return [{"content": content, "metadata": {"source": os.path.basename(path)}}]


def load_pptx(path: str) -> List[dict]:
    prs = Presentation(path)
    text_content = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        if slide_text:
            text_content.append("\n".join(slide_text))
    full_text = "\n\n".join(text_content)
    return [{"content": full_text, "metadata": {"source": os.path.basename(path)}}]


def load_html(path: str) -> List[dict]:
    docs = UnstructuredHTMLLoader(path).load()
    return [{"content": d.page_content or "", "metadata": {"source": os.path.basename(path)}} for d in docs]



def load_excel(path: str) -> List[dict]:
    import pandas as pd
    try:
        # Load all sheets
        xls = pd.ExcelFile(path)
        text_content = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            # Convert dataframe to string representation
            sheet_text = f"Sheet: {sheet_name}\n" + df.to_string(index=False)
            text_content.append(sheet_text)
        full_text = "\n\n".join(text_content)
        return [{"content": full_text, "metadata": {"source": os.path.basename(path)}}]
    except Exception as e:
        print(f"Error loading Excel file: {e}")
        return []


def load_csv(path: str) -> List[dict]:
    import pandas as pd
    try:
        df = pd.read_csv(path)
        content = "CSV Content:\n" + df.to_string(index=False)
        return [{"content": content, "metadata": {"source": os.path.basename(path)}}]
    except Exception as e:
        print(f"Error loading CSV file: {e}")
        return []

def load_any(path: str) -> List[dict]:
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        return load_pdf(path)
    if ext == ".txt":
        return load_txt(path)
    if ext == ".docx":
        return load_docx(path)
    if ext == ".pptx":
        return load_pptx(path)
    if ext in {".xlsx", ".xls"}:
        return load_excel(path)
    if ext == ".csv":
        return load_csv(path)
    if ext in {".html", ".htm"}:
        return load_html(path)

    raise ValueError(f"Unsupported file type: {ext}")
